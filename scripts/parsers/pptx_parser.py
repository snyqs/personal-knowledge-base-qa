from pptx import Presentation
from pptx.table import Table
from typing import List


class PPTXParser:
    def __init__(self):
        pass

    def parse(self, file_path: str) -> str:
        md_content = []
        try:
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                md_content.append(f"## 幻灯片{slide_num}")
                
                title = slide.shapes.title
                if title and title.text:
                    md_content.append(f"### {title.text}")
                
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text
                        if text is not None:
                            text = text.strip()
                            if text and (title is None or text != title.text):
                                md_content.append(text)
                    
                    elif shape.has_table:
                        table = shape.table
                        md_table = self._table_to_markdown(table)
                        if md_table:
                            md_content.append("\n")
                            md_content.append(md_table)
                
                md_content.append("\n---\n")
            
            return "\n".join(md_content)
        except Exception as e:
            raise RuntimeError(f"PPTX解析失败: {str(e)}")

    def _table_to_markdown(self, table: Table) -> str:
        if not table.rows:
            return ""
        
        md_lines = []
        for row_idx, row in enumerate(table.rows):
            cells = []
            for cell in row.cells:
                cell_text = cell.text_frame.text.strip()
                cell_text = cell_text.replace('\n', ' ')
                cells.append(cell_text)
            
            if row_idx == 0:
                md_lines.append(f"| {' | '.join(cells)} |")
                md_lines.append(f"| {' | '.join(['---'] * len(cells))} |")
            else:
                md_lines.append(f"| {' | '.join(cells)} |")
        
        return '\n'.join(md_lines)